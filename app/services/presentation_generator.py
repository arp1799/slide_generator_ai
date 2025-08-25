from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from typing import List
import uuid
import os
from datetime import datetime
from app.models.slide_models import SlideContent, Theme, ColorScheme, FontSettings
from app.core.config import settings


class PresentationGenerator:
    def __init__(self):
        self.themes = {
            Theme.MODERN: {
                'primary_color': RGBColor(46, 134, 171),  # #2E86AB
                'secondary_color': RGBColor(162, 59, 114),  # #A23B72
                'background_color': RGBColor(255, 255, 255),  # #FFFFFF
                'text_color': RGBColor(51, 51, 51),  # #333333
                'title_font': 'Arial',
                'body_font': 'Calibri'
            },
            Theme.CORPORATE: {
                'primary_color': RGBColor(0, 51, 102),  # #003366
                'secondary_color': RGBColor(102, 153, 204),  # #6699CC
                'background_color': RGBColor(255, 255, 255),  # #FFFFFF
                'text_color': RGBColor(0, 0, 0),  # #000000
                'title_font': 'Arial',
                'body_font': 'Calibri'
            },
            Theme.CREATIVE: {
                'primary_color': RGBColor(255, 105, 180),  # #FF69B4
                'secondary_color': RGBColor(138, 43, 226),  # #8A2BE2
                'background_color': RGBColor(255, 255, 255),  # #FFFFFF
                'text_color': RGBColor(51, 51, 51),  # #333333
                'title_font': 'Arial',
                'body_font': 'Calibri'
            },
            Theme.MINIMAL: {
                'primary_color': RGBColor(128, 128, 128),  # #808080
                'secondary_color': RGBColor(192, 192, 192),  # #C0C0C0
                'background_color': RGBColor(255, 255, 255),  # #FFFFFF
                'text_color': RGBColor(0, 0, 0),  # #000000
                'title_font': 'Arial',
                'body_font': 'Calibri'
            }
        }
    
    def generate_presentation(
        self, 
        slides: List[SlideContent], 
        theme: Theme = Theme.MODERN,
        color_scheme: ColorScheme = None,
        font_settings: FontSettings = None,
        include_citations: bool = True
    ) -> str:
        """Generate a PowerPoint presentation from slide content"""
        
        # Create presentation
        prs = Presentation()
        
        # Apply theme
        theme_colors = self.themes[theme]
        if color_scheme:
            theme_colors = self._convert_color_scheme(color_scheme)
        
        # Apply fonts
        title_font = font_settings.title_font if font_settings else theme_colors['title_font']
        body_font = font_settings.body_font if font_settings else theme_colors['body_font']
        title_size = font_settings.title_size if font_settings else 44
        body_size = font_settings.body_size if font_settings else 18
        
        # Generate slides
        for i, slide_content in enumerate(slides):
            if i == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(slide_layout)
                self._create_title_slide(slide, slide_content, theme_colors, title_font, title_size)
            else:
                # Content slide
                slide_layout = prs.slide_layouts[1]  # Content slide layout
                slide = prs.slides.add_slide(slide_layout)
                self._create_content_slide(slide, slide_content, theme_colors, title_font, body_font, title_size, body_size)
        
        # Add citations slide if requested
        if include_citations and len(slides) > 1:
            self._add_citations_slide(prs, theme_colors, title_font, body_font, title_size, body_size)
        
        # Generate filename and save
        filename = self._generate_filename(slides[0].title if slides else "presentation")
        filepath = os.path.join(settings.output_dir, filename)
        prs.save(filepath)
        
        return filename
    
    def _create_title_slide(self, slide, slide_content: SlideContent, theme_colors: dict, title_font: str, title_size: int):
        """Create a title slide"""
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme_colors['background_color']
        
        # Add title
        title = slide.shapes.title
        title.text = slide_content.title
        title.text_frame.paragraphs[0].font.name = title_font
        title.text_frame.paragraphs[0].font.size = Pt(title_size)
        title.text_frame.paragraphs[0].font.color.rgb = theme_colors['primary_color']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add subtitle if content exists
        if slide_content.content:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_content.content
            subtitle.text_frame.paragraphs[0].font.name = title_font
            subtitle.text_frame.paragraphs[0].font.size = Pt(title_size // 2)
            subtitle.text_frame.paragraphs[0].font.color.rgb = theme_colors['secondary_color']
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, slide, slide_content: SlideContent, theme_colors: dict, title_font: str, body_font: str, title_size: int, body_size: int):
        """Create a content slide based on layout"""
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme_colors['background_color']
        
        # Add title
        title = slide.shapes.title
        title.text = slide_content.title
        title.text_frame.paragraphs[0].font.name = title_font
        title.text_frame.paragraphs[0].font.size = Pt(title_size)
        title.text_frame.paragraphs[0].font.color.rgb = theme_colors['primary_color']
        
        # Create content based on layout
        if slide_content.layout.value == "bullet_points":
            self._create_bullet_points_slide(slide, slide_content, theme_colors, body_font, body_size)
        elif slide_content.layout.value == "two_column":
            self._create_two_column_slide(slide, slide_content, theme_colors, body_font, body_size)
        elif slide_content.layout.value == "content_with_image":
            self._create_content_with_image_slide(slide, slide_content, theme_colors, body_font, body_size)
        else:
            # Default to bullet points
            self._create_bullet_points_slide(slide, slide_content, theme_colors, body_font, body_size)
    
    def _create_bullet_points_slide(self, slide, slide_content: SlideContent, theme_colors: dict, body_font: str, body_size: int):
        """Create a bullet points slide"""
        # Clear existing content
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.text_frame.clear()
        
        # Add bullet points
        if slide_content.bullet_points:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            text_frame = textbox.text_frame
            text_frame.clear()
            
            for i, point in enumerate(slide_content.bullet_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = f"• {point}"
                p.font.name = body_font
                p.font.size = Pt(body_size)
                p.font.color.rgb = theme_colors['text_color']
                p.level = 0
        
        # Add content if no bullet points
        elif slide_content.content:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            text_frame = textbox.text_frame
            text_frame.text = slide_content.content
            text_frame.paragraphs[0].font.name = body_font
            text_frame.paragraphs[0].font.size = Pt(body_size)
            text_frame.paragraphs[0].font.color.rgb = theme_colors['text_color']
    
    def _create_two_column_slide(self, slide, slide_content: SlideContent, theme_colors: dict, body_font: str, body_size: int):
        """Create a two-column slide"""
        # Clear existing content
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.text_frame.clear()
        
        # Left column
        if slide_content.left_column:
            left_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5))
            left_frame = left_textbox.text_frame
            left_frame.text = slide_content.left_column
            left_frame.paragraphs[0].font.name = body_font
            left_frame.paragraphs[0].font.size = Pt(body_size)
            left_frame.paragraphs[0].font.color.rgb = theme_colors['text_color']
        
        # Right column
        if slide_content.right_column:
            right_textbox = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(5))
            right_frame = right_textbox.text_frame
            right_frame.text = slide_content.right_column
            right_frame.paragraphs[0].font.name = body_font
            right_frame.paragraphs[0].font.size = Pt(body_size)
            right_frame.paragraphs[0].font.color.rgb = theme_colors['text_color']
    
    def _create_content_with_image_slide(self, slide, slide_content: SlideContent, theme_colors: dict, body_font: str, body_size: int):
        """Create a content slide with image placeholder"""
        # Clear existing content
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.text_frame.clear()
        
        # Add content
        if slide_content.content:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(4))
            text_frame = textbox.text_frame
            text_frame.text = slide_content.content
            text_frame.paragraphs[0].font.name = body_font
            text_frame.paragraphs[0].font.size = Pt(body_size)
            text_frame.paragraphs[0].font.color.rgb = theme_colors['text_color']
        
        # Add image placeholder
        if slide_content.image_placeholder:
            # Create a placeholder rectangle
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.5), Inches(2), Inches(3), Inches(4)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = theme_colors['secondary_color']
            
            # Add placeholder text
            placeholder_textbox = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(3), Inches(4))
            placeholder_frame = placeholder_textbox.text_frame
            placeholder_frame.text = f"[Image Placeholder]\n{slide_content.image_placeholder}"
            placeholder_frame.paragraphs[0].font.name = body_font
            placeholder_frame.paragraphs[0].font.size = Pt(body_size - 2)
            placeholder_frame.paragraphs[0].font.color.rgb = theme_colors['background_color']
            placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_citations_slide(self, prs, theme_colors: dict, title_font: str, body_font: str, title_size: int, body_size: int):
        """Add a citations slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme_colors['background_color']
        
        # Add title
        title = slide.shapes.title
        title.text = "References & Citations"
        title.text_frame.paragraphs[0].font.name = title_font
        title.text_frame.paragraphs[0].font.size = Pt(title_size)
        title.text_frame.paragraphs[0].font.color.rgb = theme_colors['primary_color']
        
        # Add citations
        citations = [
            "• Research and content generated using AI technology",
            "• Presentation created with Slide Generator API",
            "• For educational and demonstration purposes",
            "• Please verify all information before use in professional settings"
        ]
        
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        text_frame = textbox.text_frame
        text_frame.clear()
        
        for i, citation in enumerate(citations):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = citation
            p.font.name = body_font
            p.font.size = Pt(body_size - 2)
            p.font.color.rgb = theme_colors['text_color']
    
    def _convert_color_scheme(self, color_scheme: ColorScheme) -> dict:
        """Convert ColorScheme to RGB values"""
        def hex_to_rgb(hex_color: str) -> RGBColor:
            hex_color = hex_color.lstrip('#')
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
        
        return {
            'primary_color': hex_to_rgb(color_scheme.primary_color),
            'secondary_color': hex_to_rgb(color_scheme.secondary_color),
            'background_color': hex_to_rgb(color_scheme.background_color),
            'text_color': hex_to_rgb(color_scheme.text_color),
            'title_font': 'Arial',
            'body_font': 'Calibri'
        }
    
    def _generate_filename(self, title: str) -> str:
        """Generate a unique filename for the presentation"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        safe_title = safe_title.replace(' ', '_')[:30]
        unique_id = str(uuid.uuid4())[:8]
        return f"{safe_title}_{timestamp}_{unique_id}.pptx" 